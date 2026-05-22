"""
ingest_corpus.py — Ingest filesystem corpus to pgvector with module/week metadata.

Walks `corpus/semana-XX/sesion-YY/{antes,durante,despues}/*.{pdf,docx,pptx,txt,md}`
and `corpus/recursos-kotlin/*`, parses content, chunks via RecursiveCharacterTextSplitter,
embeds with mxbai-embed-large and stores in `document_chunks` with rich metadata.

Each file becomes one `documents` row (status=active) so admin UI lists them.
Each chunk's `metadata` jsonb carries: {module_id, semana, sesion, fase, filename, source_path}.

Usage (recommended — runs inside backend container with corpus mounted ro):
    docker exec tutor_backend python scripts/ingest_corpus.py
    docker exec tutor_backend python scripts/ingest_corpus.py --reset   # drop previous corpus docs and re-ingest

Module mapping derived from docs/corpus-mapping.md (semana → M1..M5):
    01-02 → M1 Fundamentos
    03-08 → M2 Kotlin
    09-10 → M3 UI
    11-13 → M4 Componentes/Datos
    14-16 → M5 Avanzado/Despliegue
    recursos-kotlin → M2 (refuerzo)
"""

from __future__ import annotations

import argparse
import asyncio
import mimetypes
import os
import sys
import uuid
from datetime import datetime, timezone
from pathlib import Path

sys.path.insert(0, os.path.join(os.path.dirname(__file__), ".."))

from sqlalchemy import select, delete, text
from sqlalchemy.ext.asyncio import create_async_engine, AsyncSession
from sqlalchemy.orm import sessionmaker

from app.config import settings
from app.models.document import Document, DocumentChunk
from app.utils.chunking import get_text_splitter

from langchain_ollama import OllamaEmbeddings
from pypdf import PdfReader
from docx import Document as DocxDocument

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False


SEMANA_TO_MODULE: dict[str, int] = {
    "01": 1, "02": 1,
    "03": 2, "04": 2, "05": 2, "06": 2, "07": 2, "08": 2,
    "09": 3, "10": 3,
    "11": 4, "12": 4, "13": 4,
    "14": 5, "15": 5, "16": 5,
}

# Container path when running via docker exec (corpus mounted ro at /app/corpus)
# Falls back to repo root if running from host.
CORPUS_CANDIDATES = [
    Path("/app/corpus"),
    Path(__file__).resolve().parents[2] / "corpus",
]
CORPUS_ROOT = next((p for p in CORPUS_CANDIDATES if p.exists()), CORPUS_CANDIDATES[0])

BATCH_SIZE = 5
MIN_TEXT_LEN = 100
SKIP_DIRS = {"_raw"}
SKIP_FILENAMES = {
    "antes.txt",
    "despues.txt",
    "Nuevo Documento de texto.txt",
    "Nuevo Documento de Microsoft Word.docx",
}
SUPPORTED_EXTS = {".pdf", ".docx", ".pptx", ".txt", ".md"}
SOURCE_PREFIX = "corpus://"


def extract_text(path: Path) -> str:
    ext = path.suffix.lower()
    try:
        if ext == ".pdf":
            reader = PdfReader(str(path))
            return "\n".join((p.extract_text() or "") for p in reader.pages)
        if ext == ".docx":
            doc = DocxDocument(str(path))
            paragraphs = [p.text for p in doc.paragraphs if p.text]
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        if cell.text:
                            paragraphs.append(cell.text)
            return "\n".join(paragraphs)
        if ext == ".pptx":
            if not HAS_PPTX:
                return ""
            prs = Presentation(str(path))
            pieces: list[str] = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        pieces.append(shape.text)
            return "\n".join(pieces)
        if ext in {".txt", ".md"}:
            return path.read_text(encoding="utf-8", errors="ignore")
    except Exception as exc:
        print(f"  WARN extract error {path.name}: {exc}")
    return ""


def parse_path(path: Path) -> dict:
    rel = path.relative_to(CORPUS_ROOT)
    parts = rel.parts
    meta: dict = {
        "source_path": str(rel).replace("\\", "/"),
        "filename": path.name,
        "module_id": None,
        "semana": None,
        "sesion": None,
        "fase": None,
    }
    for part in parts:
        if part.startswith("semana-"):
            num = part.replace("semana-", "")[:2]
            meta["semana"] = num
            meta["module_id"] = SEMANA_TO_MODULE.get(num)
        elif part.startswith("sesion-"):
            meta["sesion"] = part.replace("sesion-", "")
        elif part in {"antes", "durante", "despues"}:
            meta["fase"] = part
    if parts and parts[0] == "recursos-kotlin":
        meta["module_id"] = 2
        meta["fase"] = "recurso"
    return meta


def find_candidates() -> list[Path]:
    out: list[Path] = []
    for path in CORPUS_ROOT.rglob("*"):
        if not path.is_file():
            continue
        if any(skip in path.parts for skip in SKIP_DIRS):
            continue
        if path.name in SKIP_FILENAMES:
            continue
        if path.suffix.lower() not in SUPPORTED_EXTS:
            continue
        out.append(path)
    out.sort()
    return out


async def main(reset: bool = False) -> None:
    print("=" * 60)
    print("INGESTA CORPUS FILESYSTEM -> pgvector (módulos M1-M5)")
    print("=" * 60)
    print(f"Corpus root: {CORPUS_ROOT}")
    if not CORPUS_ROOT.exists():
        print(f"ERROR: corpus root no existe: {CORPUS_ROOT}")
        return
    if not HAS_PPTX:
        print("WARN: python-pptx no instalado — archivos .pptx serán omitidos.")

    print(f"Ollama: {settings.OLLAMA_BASE_URL} model={settings.OLLAMA_EMBED_MODEL}")
    embeddings = OllamaEmbeddings(
        model=settings.OLLAMA_EMBED_MODEL,
        base_url=settings.OLLAMA_BASE_URL,
    )
    try:
        dim = len(await embeddings.aembed_query("test"))
        print(f"OK embeddings dim={dim}")
    except Exception as exc:
        print(f"ERROR Ollama no responde: {exc}")
        return

    splitter = get_text_splitter()
    engine = create_async_engine(settings.DATABASE_URL)
    SessionLocal = sessionmaker(engine, class_=AsyncSession, expire_on_commit=False)

    candidates = find_candidates()
    print(f"\nArchivos candidatos: {len(candidates)}")

    async with SessionLocal() as db:
        if reset:
            print("\nReset: borrando documentos previos corpus://...")
            doc_ids = (await db.execute(
                select(Document.id).where(Document.original_filename.like(f"{SOURCE_PREFIX}%"))
            )).scalars().all()
            if doc_ids:
                await db.execute(delete(DocumentChunk).where(DocumentChunk.document_id.in_(doc_ids)))
                await db.execute(delete(Document).where(Document.id.in_(doc_ids)))
                await db.commit()
                print(f"  borrados {len(doc_ids)} documentos previos")

        stats = {"ingested": 0, "skipped_existing": 0, "skipped_short": 0,
                 "skipped_parse": 0, "errors": 0, "chunks": 0}
        per_module: dict[int, int] = {}

        for idx, path in enumerate(candidates, 1):
            meta_path = parse_path(path)
            source_key = f"{SOURCE_PREFIX}{meta_path['source_path']}"
            tag = f"[{idx}/{len(candidates)}] M{meta_path.get('module_id') or '?'}"

            if not reset:
                existing = (await db.execute(
                    select(Document).where(Document.original_filename == source_key)
                )).scalar_one_or_none()
                if existing and existing.status == "active":
                    stats["skipped_existing"] += 1
                    continue

            raw = extract_text(path)
            if len(raw.strip()) < MIN_TEXT_LEN:
                stats["skipped_short"] += 1
                print(f"{tag} SKIP-short {path.name}")
                continue

            chunks = splitter.split_text(raw)
            if not chunks:
                stats["skipped_parse"] += 1
                print(f"{tag} SKIP-empty {path.name}")
                continue

            doc = Document(
                id=uuid.uuid4(),
                original_filename=source_key,
                stored_filename=path.name[:500],
                file_size_bytes=path.stat().st_size,
                mime_type=mimetypes.guess_type(path.name)[0] or "application/octet-stream",
                status="processing",
                uploaded_by=None,
            )
            db.add(doc)
            await db.commit()
            await db.refresh(doc)

            embeds: list[list[float] | None] = []
            for i in range(0, len(chunks), BATCH_SIZE):
                batch = chunks[i:i + BATCH_SIZE]
                try:
                    batch_emb = await embeddings.aembed_documents(batch)
                    embeds.extend(batch_emb)
                except Exception as exc:
                    print(f"  WARN embed batch error: {exc}")
                    embeds.extend([None] * len(batch))

            saved = 0
            for i, (ctext, emb) in enumerate(zip(chunks, embeds)):
                if emb is None:
                    continue
                cmeta = {
                    **meta_path,
                    "source": f"M{meta_path.get('module_id') or '?'} > "
                              f"{meta_path.get('semana') or '-'}/"
                              f"{meta_path.get('sesion') or '-'} > {path.name}",
                }
                db.add(DocumentChunk(
                    id=uuid.uuid4(),
                    document_id=doc.id,
                    content=ctext,
                    embedding=emb,
                    chunk_index=i,
                    metadata_=cmeta,
                ))
                saved += 1

            if saved == 0:
                doc.status = "error"
                doc.error_message = "embeddings failed for all chunks"
                stats["errors"] += 1
            else:
                doc.status = "active"
                doc.chunk_count = saved
                doc.processed_at = datetime.now(timezone.utc)
                stats["ingested"] += 1
                stats["chunks"] += saved
                mid = meta_path.get("module_id") or 0
                per_module[mid] = per_module.get(mid, 0) + saved
            await db.commit()
            print(f"{tag} OK {path.name} -> {saved} chunks")

        print("\n" + "=" * 60)
        print("RESUMEN")
        print("=" * 60)
        print(f"Documentos ingestados : {stats['ingested']}")
        print(f"Total chunks          : {stats['chunks']}")
        print(f"Saltados (ya existían): {stats['skipped_existing']}")
        print(f"Saltados (texto corto): {stats['skipped_short']}")
        print(f"Saltados (parse vacío): {stats['skipped_parse']}")
        print(f"Errores               : {stats['errors']}")
        print("\nChunks por módulo:")
        for mid in sorted(per_module):
            label = f"M{mid}" if mid else "(sin módulo)"
            print(f"  {label}: {per_module[mid]}")

        print("\nVerificación búsqueda semántica...")
        try:
            q = "¿Cómo declarar variables en Kotlin?"
            qvec = await embeddings.aembed_query(q)
            vlit = "[" + ",".join(f"{x}" for x in qvec) + "]"
            # asyncpg + pgvector: vector literal must be inlined (parameter ::vector cast fails).
            sql = f"""
                SELECT content, metadata AS meta,
                       1 - (embedding <=> '{vlit}'::vector) AS similarity
                FROM document_chunks
                WHERE document_id IN (
                    SELECT id FROM documents WHERE original_filename LIKE :prefix
                )
                ORDER BY embedding <=> '{vlit}'::vector
                LIMIT 3
            """
            result = await db.execute(text(sql), {"prefix": f"{SOURCE_PREFIX}%"})
            for j, row in enumerate(result.fetchall(), 1):
                src = (row.meta or {}).get("source", "?")
                print(f"  [{j}] sim={row.similarity:.3f}  {src}")
                print(f"      {row.content[:120]}...")
        except Exception as exc:
            print(f"  WARN verificación: {exc}")

    await engine.dispose()


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Ingest corpus filesystem → pgvector")
    parser.add_argument("--reset", action="store_true",
                        help="Borra documentos corpus:// previos y re-ingesta")
    args = parser.parse_args()
    asyncio.run(main(reset=args.reset))
